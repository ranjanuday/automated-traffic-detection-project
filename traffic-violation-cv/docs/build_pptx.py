"""Generate the AutoTraffic Vision PowerPoint deck (.pptx).

Reproducible: `python docs/build_pptx.py` -> docs/AutoTraffic_Vision.pptx
Mirrors the content of docs/slides.html in native PowerPoint form.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --- Palette (neutral, no brand references) --------------------------------
NAVY = RGBColor(0x0D, 0x24, 0x40)
BLUE = RGBColor(0x15, 0x65, 0xC0)
AMBER = RGBColor(0xFF, 0xC2, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF3, 0xF5, 0xF8)

TITLE = "AutoTraffic Vision"
SUBTITLE = "AI-Powered Traffic Violation Detection from Photographic Evidence"
DESCRIPTION = (
    "A scalable computer-vision system that automatically processes traffic "
    "surveillance images, detects vehicles and road users, identifies and "
    "classifies seven types of traffic violations, recognizes license plates "
    "via OCR, generates annotated evidence with metadata and timestamps, and "
    "provides analytics, searchable records, and performance evaluation - "
    "drastically reducing manual inspection effort for traffic enforcement."
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def _text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space=6):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def _bar(slide):
    """Top accent bar for content slides."""
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()


def content_slide(title, bullets, footer=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _bar(s)
    tb = _box(s, 0.7, 0.5, 12.0, 1.0)
    _text(tb.text_frame, title, 30, NAVY, bold=True)
    body = _box(s, 0.9, 1.7, 11.6, 5.2)
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        lvl = 0
        txt = b
        if isinstance(b, tuple):
            txt, lvl = b
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = ("- " if lvl else "") + txt
        r.font.size = Pt(20 if lvl == 0 else 16)
        r.font.bold = lvl == 0
        r.font.color.rgb = NAVY if lvl == 0 else GRAY
        r.font.name = "Segoe UI"
    if footer:
        fb = _box(s, 0.9, 6.9, 11.6, 0.5)
        _text(fb.text_frame, footer, 12, GRAY)
    return s


# --- Slide 1: Title --------------------------------------------------------
s = prs.slides.add_slide(BLANK)
_bg(s, NAVY)
band = s.shapes.add_shape(1, 0, Inches(2.9), SW, Inches(0.08))
band.fill.solid(); band.fill.fore_color.rgb = AMBER; band.line.fill.background()
t = _box(s, 0.8, 1.5, 11.7, 1.3)
_text(t.text_frame, TITLE, 54, WHITE, bold=True, align=PP_ALIGN.CENTER)
st = _box(s, 1.0, 3.1, 11.3, 1.0)
_text(st.text_frame, SUBTITLE, 24, AMBER, align=PP_ALIGN.CENTER)
d = _box(s, 1.6, 4.3, 10.1, 2.0)
_text(d.text_frame, DESCRIPTION, 15, RGBColor(0xCF, 0xD8, 0xE3),
      align=PP_ALIGN.CENTER)
tag = _box(s, 1.0, 6.6, 11.3, 0.5)
_text(tag.text_frame, "Computer Vision for Smart Traffic Enforcement", 14,
      RGBColor(0x9A, 0xA7, 0xB8), align=PP_ALIGN.CENTER)

# --- Slide 2: Problem ------------------------------------------------------
content_slide("The Problem", [
    "Traffic surveillance cameras generate huge volumes of images every day.",
    ("Manual inspection is labor-intensive and slow", 1),
    ("Human review is inconsistent and error-prone", 1),
    ("Enforcement does not scale with rising traffic volume", 1),
    "Need: an intelligent system to auto-analyze photographic evidence and "
    "improve the efficiency and accuracy of traffic-law enforcement.",
])

# --- Slide 3: Solution Overview -------------------------------------------
content_slide("Solution Overview", [
    "An end-to-end computer-vision pipeline:",
    ("Preprocess -> Detect -> Classify Violations -> Read Plate -> "
     "Generate Evidence -> Analytics -> Evaluate", 1),
    "Robust across challenging conditions:",
    ("Low light, rain, shadows, motion blur", 1),
    ("Varying traffic density and image quality", 1),
])

# --- Slide 4: Architecture -------------------------------------------------
content_slide("System Architecture", [
    "Pipeline stages (left to right):",
    ("Image input", 1),
    ("Preprocessing: CLAHE + gamma + denoise (optional deep Restormer "
     "de-rain / de-blur)", 1),
    ("Object detection: COCO YOLOv8 + Indian-vehicle YOLO fusion", 1),
    ("7 violation detectors (rule + model based)", 1),
    ("License plate: YOLOv5 detector + EasyOCR", 1),
    ("Evidence: annotated image + SQLite metadata + timestamp", 1),
    ("Analytics dashboard", 1),
], footer="Tech stack: Python, FastAPI, HTMX, Tailwind, Chart.js, SQLite, "
          "Ultralytics YOLO, EasyOCR, Restormer")

# --- Slide 5: Preprocessing ------------------------------------------------
content_slide("Image Preprocessing", [
    "Adaptive enhancement keeps detection robust across conditions:",
    ("Low light: gamma correction + CLAHE + denoise", 1),
    ("Rain / haze: CLAHE + bilateral filtering", 1),
    ("Shadows and motion blur handling", 1),
    "Optional deep restoration with Restormer (de-rain, de-blur) models for "
    "heavy-weather robustness.",
])

# --- Slide 6: Detection ----------------------------------------------------
content_slide("Vehicle & Road-User Detection", [
    "Sensor fusion of two detectors:",
    ("COCO YOLOv8: person, car, bus, truck, motorcycle, bicycle, "
     "traffic light", 1),
    ("South-Asian / Indian-vehicle YOLO: adds auto-rickshaw and rickshaw, "
     "re-labels three-wheeler mislabels", 1),
    "Correctly localizes and classifies vehicle categories and road users.",
])

# --- Slide 7: 7 Violations -------------------------------------------------
content_slide("The 7 Violation Types", [
    ("No Helmet - YOLOv8 helmet/head model (proven on real riders)", 1),
    ("Triple Riding - 3+ persons on one two-wheeler (geometry)", 1),
    ("No Seatbelt - YOLOv11 classifier (windscreen-camera opt-in)", 1),
    ("Red-Light Jump - HSV signal state + stop-line crossing", 1),
    ("Stop-Line Violation - configured or auto-detected white marking", 1),
    ("Wrong-Side Driving - configured divider (advisory)", 1),
    ("Illegal Parking - configured no-parking zones", 1),
    "Every prediction carries a confidence score.",
])

# --- Slide 8: ANPR ---------------------------------------------------------
content_slide("License Plate Recognition", [
    "Two-stage Automatic Number-Plate Recognition (ANPR):",
    ("Stage 1: YOLOv5 plate detector tightly crops the plate", 1),
    ("Stage 2: EasyOCR extracts the registration string", 1),
    "Graceful degradation: returns nothing rather than garbage on unreadable "
    "plates.",
])

# --- Slide 9: Evidence -----------------------------------------------------
content_slide("Evidence Generation", [
    "Annotated images with colored bounding boxes + violation labels + "
    "confidence (e.g. \"No Helmet 73%\").",
    "Metadata persisted to SQLite:",
    ("Violation type, confidence, bounding box, plate, note", 1),
    ("UTC timestamps for every record", 1),
])

# --- Slide 10: Analytics ---------------------------------------------------
content_slide("Analytics & Reporting", [
    "Dashboard built with HTMX + Chart.js:",
    ("Violation statistics by type", 1),
    ("Daily / trend breakdowns", 1),
    ("Top offending plates", 1),
    "Searchable records filterable by violation type and plate number.",
])

# --- Slide 11: Performance -------------------------------------------------
content_slide("Performance Evaluation", [
    "Metrics harness computes:",
    ("Accuracy, Precision, Recall, F1-score per class", 1),
    ("Real mAP@0.5 via IoU box matching", 1),
    ("Per-stage computational efficiency (latency)", 1),
    "On the labelled sample set: 100% image-level accuracy, no_helmet "
    "F1 = 1.00; ~6.5 s/image on CPU.",
])

# --- Slide 12: Engineering Quality ----------------------------------------
content_slide("Engineering Quality", [
    "Real pretrained models (not mocks); graceful degradation.",
    "False-positive guards proven by the eval harness:",
    ("A red parked car is NOT mistaken for a red signal", 1),
    ("A sunlit road is NOT mistaken for a stop line", 1),
    "Detectors abstain when un-configured rather than fabricate violations.",
])

# --- Slide 13: Limitations -------------------------------------------------
content_slide("Limitations & Future Work", [
    "Honest current limits:",
    ("Seatbelt model imperfect on out-of-distribution images", 1),
    ("Wrong-side from a still frame is advisory (needs video/motion)", 1),
    ("Small labelled eval set; needs hundreds of annotated images", 1),
    ("Config-driven detectors need per-camera setup", 1),
    "Future: video tracking, larger detectors, region-fine-tuned models, "
    "scalable deployment.",
])

# --- Slide 14: Summary -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
_bg(s, NAVY)
band = s.shapes.add_shape(1, 0, Inches(2.0), SW, Inches(0.08))
band.fill.solid(); band.fill.fore_color.rgb = AMBER; band.line.fill.background()
h = _box(s, 0.8, 1.0, 11.7, 1.0)
_text(h.text_frame, "Summary & Impact", 40, WHITE, bold=True,
      align=PP_ALIGN.CENTER)
b = _box(s, 1.4, 2.5, 10.5, 3.5)
_text(b.text_frame,
      "A working, scalable AI traffic-image analysis system that "
      "automatically identifies, classifies, and documents 7 violation types "
      "from photographic evidence - reducing manual effort and improving the "
      "effectiveness of traffic monitoring and enforcement.",
      20, RGBColor(0xDD, 0xE5, 0xEF), align=PP_ALIGN.CENTER)

out = Path(__file__).parent / "AutoTraffic_Vision.pptx"
prs.save(str(out))
print(f"Saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
